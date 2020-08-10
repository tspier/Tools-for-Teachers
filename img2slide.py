from pptx import Presentation 
from pptx.util import Inches

#all_image_names = ['001', '002', '003']
all_image_names = []

prs = Presentation() 
blank_slide_layout = prs.slide_layouts[6]

def add_slide(img_path):
    slide = prs.slides.add_slide(blank_slide_layout)

    # This will force the image to be placed in the
    # top-left corner, so no spacing will be added.
    left = Inches(0.0)
    top = Inches(0.0)
    
    # This assigns the file names to img_path and .jpg, e.g. "001" + .jpg = 001.jpg
    # Change this if you are going to specify the extension names OR have images
    # of different file extensions. You should also specify the width and height for
    # the images to be included. Use an online ratio calculator if necessary to ensure
    # that the images are not distorted. It looks like the total available size is 12.63
    # for height and 20.05 for width.
    pic = slide.shapes.add_picture(img_path + ".jpg", left, top, width=Inches(0), height=Inches(0))

print("Creating Presentation...")

for image in all_image_names:
    add_slide(image)
    print("Adding Slide #" + str(image) + "...")

print("Finished.")
prs.save('output.pptx')
