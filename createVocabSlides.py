from bs4 import BeautifulSoup
from pptx import Presentation
import time
import requests

def findEtymology(query):
    url = f'https://www.dictionary.com/browse/{query}?s=t'

    r = requests.get(url)
    soup = BeautifulSoup(r.content, 'html.parser')
    wordDefinition = soup.find('span', attrs={'one-click-content css-1p89gle e1q3nk1v4'})
    wordEtymology = soup.find('div', attrs={'one-click-content css-otpbu9 e16svm7n0'})

    print(query)
    print(wordDefinition.text)
    print(wordEtymology.text)
    print("\n")

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = query

    tf = body_shape.text_frame
    tf.text = wordEtymology.text + "\n" + wordDefinition.text


# allWords = "zealot derogatory omnipotent contingent stringent"
allWords = ""

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

for x in allWords.split(" "):
    findEtymology(x)

prs.save('output.pptx')
