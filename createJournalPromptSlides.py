from pptx import Presentation

"""
quotes = {
    "THE YOUNG QUEEN:": "Write about a character who survived the apocalypse and has now taken the role of leader for a small group of survivors. The catch? This character is the oldest person in the group. Double catch? She is only seventeen-years-old.",
    "MASS AMNESIA:": "Write about a character who lives in a city that experienced mass amnesia one year ago. The catch? No one remembers how or why it happened. Double catch? This character's memory was not at all affected by the mass amnesia."
    }
"""
quotes = {}
	
def addSlide(quote, person):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    #title_shape.text = "Journal Prompt"
    title_shape.text = ""

    tf = body_shape.text_frame
    tf.text = quote + "\n" + person

prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

for x in quotes:
    addSlide(x, quotes[x])

prs.save('output.pptx')
